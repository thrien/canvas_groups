#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Utilities for using Canvas as a GSI

It is intented to be used by GSIs for PHYSICS 151/251 at the University of
Michigan.
"""
# standard library
import os.path
import argparse
from argparse import ArgumentTypeError
from urllib.request import Request, urlopen
import json
# external dependencies
# for sign-in sheets
import numpy as np
import matplotlib.pyplot as plt
# for introduction slides
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# configure help message formatting
class RawDescriptionDefaultsHelpFormatter(
        argparse.RawDescriptionHelpFormatter,
        argparse.ArgumentDefaultsHelpFormatter):
    pass

# Canvas API
API_URL = "https://umich.instructure.com/api/v1"
TOKEN = ""
COURSE_ID = 850281

# Make sure we are in the right directory
os.chdir(os.path.dirname(__file__))

# already existing labs on disk
prefix = "lab"
existing_labs = [int(entry.removeprefix(prefix)) for entry in os.listdir()
                 if os.path.isdir(entry) and entry.startswith(prefix)]
existing_labs.sort()

instructor = ""  # last name, first name
# int for group numbers, 'I' for the instructor (optional), '.' for nothing
table_layout = [[ 1 , 'I', '.'],
                [ 2 , '.',  8 ],
                [ 3 , '.',  7 ],
                [ 4 ,  5 ,  6 ]]
tables = [table for row in table_layout
                for table in row
                if isinstance(table, int)]


def _canvas_api(command, full_url=False, headers={}, verbose=False):
    """Call the Canvas API with the given command and parameters
    
    Args:
        command: the API command, e.g., "courses/123/groups"
        parameters: additional parameters to be sent as headers, e.g.,
                    {"search_term": "lab 1"}
    Returns:
        the response as a Python object (list for JSON, else string)
    """
    if not TOKEN:
        raise RuntimeError("No Canvas API access token defined.")
    
    request = Request(command if full_url else f"{API_URL}/{command}",
                      headers={"Authorization": f"Bearer {TOKEN}"}|headers)
    response = urlopen(request)
    if response.getheader("Content-Type").startswith("application/json"):
        content = json.load(response)
    else:
        content = response.read()

    try:  # reading next pages
        links = response.getheader("Link").split(",")
        pages = {rel.removeprefix(" rel=").strip('"'): link.strip("<>")
                 for link, rel in map(lambda page: page.split(";"), links)}
        content += _canvas_api(pages["next"], full_url=True, verbose=verbose)
    except (KeyError, AttributeError):
        pass

    return content


def _canvas_import_csv(lab, verbose=False):
    # The groups for each lab are defined in a group category on Canvas
    categories = _canvas_api(f"courses/{COURSE_ID}/group_categories",
                             verbose=verbose)
    try:
        category = next(category for category in categories
                        if category["name"].lower() == f"lab {lab:d}")
    except StopIteration:
        raise RuntimeError(f"Couldn't find groups for lab {lab:d} on Canvas.")

    # Ask Canvas to export the groups for this lab as a CSV
    data = _canvas_api(f"group_categories/{category["id"]}/export",
                       verbose=verbose)

    # Save CSV on disk
    filename = os.path.join(f"lab{lab:02d}", "canvas.csv")
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    with open(filename, "wb") as file:
        file.write(data)
    if verbose:
        print(f'CSV file written to "{filename}".')


def _format_name(name):
    """Format "Smith, Emma Marie" as "Emma Smith"
    """
    lasts, firsts = name.split(",")
    first = firsts.strip().split(" ")[0]
    last = lasts.strip()  #.split(" ")[0]

    return f"__ {first:s} {last:s}"


def _draw(names, groups, title="Groups", smallfont=18, bigfont=25,
          margin=0.015, title_margin=0.075):
    """Draw the group assignment on tables

    Args:
        names: list of the names of all students
        groups: list of the corresponding group numbers
    Returns:
        the created matplotlib figure
    """
    # Create a grid of subplots based on table layout
    fig, axes = plt.subplot_mosaic(table_layout, figsize=(11, 8.5))

    fig.suptitle(title, fontsize=bigfont)

    for table in tables:
        # Remove ticks and thicken spine for group axes
        ax = axes[table]
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_linewidth(2)
        # group number in the top left corner
        ax.text(0.05, 0.90, "{:d}".format(table),
                transform=ax.transAxes,
                fontsize=smallfont, ha="left", va="top")
        # the list of names
        ax.text(0.05, 0.70, "\n".join(sorted(names[groups == table])),
                transform=ax.transAxes,
                fontsize=smallfont, ha="left", va="top")

    # Subplot for the instructor (me)
    if 'I' in axes.keys():
        ax = axes['I']
        ax.set_axis_off()
        if instructor:
            ax.text(0.5, 0.5, _format_name(instructor),
                    transform=ax.transAxes,
                    fontsize=smallfont, ha="center", va="center")

    # Adjust spacing between subplots
    fig.subplots_adjust(left=margin, right=1 - margin,
                        top=1 - title_margin, bottom=margin,
                        wspace=margin * len(table_layout[0]),
                        hspace=margin * len(table_layout))

    return fig


class FlatListAction(argparse.Action):
    """Action for argparse that flattens a list of lists

    With the interval type this parses "-l 1..3 5" into [1, 2, 3, 5].
    """
    def __call__(self, parser, namespace, values, option_string=None):
        setattr(namespace, self.dest, sum(values, start=[]))


def _interval(string, last=max(existing_labs, default=0)):
    """integers or intervals like a..b"""
    wraps = lambda s: last + 1 - int(s[1:]) if s.startswith("-") else int(s)
    if ".." in string:
        a, b = map(wraps, string.split(".."))
        return list(range(a, b + 1))
    else:
        return [wraps(string)]


def _sheets_parser(parser):
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="print status messages")
    parser.add_argument("-f", "--force", action="store_true",
                        help="pull recent CSV from Canvas")
    parser.add_argument("-e", "--extensions", nargs="+",
                        default=["pdf", "png"],
                        metavar="ext", help="output formats")
    parser.add_argument("-l", "--labs", type=_interval, nargs="+",
                        action=FlatListAction,
                        default=[max(existing_labs, default=0) + bool(TOKEN)]
                                if existing_labs or TOKEN else [],
                        metavar="numbers", help=_interval.__doc__)
    parser.add_argument("-s", "--sections", type=int, nargs="+",
                        default=[15, 25], metavar="section",
                        help="your section numbers")


def sheets(labs, sections,
           extensions=["pdf", "png"], force=False, verbose=False):
    """Draw a sign-in sheet showing what group/table students are assigned to.

    Input and output files are organized in the current directory like this:
        .
        ├── draw_sheets.py
        └── lab01
            ├── canvas.csv
            ├── groups015.png
            └── groups025.png
    """
    if verbose:
        print(f"Processing labs [{', '.join(str(lab) for lab in labs)}] and "
              f"sections [{', '.join(str(section) for section in sections)}].")

    for lab in labs:
        # Download CSV if necessary
        canvas_file = os.path.join(f"lab{lab:02d}", "canvas.csv")
        if force or not os.path.exists(canvas_file):
            if verbose:
                print(f"Downloading lab {lab:d} from Canvas.")
            # TODO catch network error to simplify error message
            _canvas_import_csv(lab, verbose=verbose)
        else:
            if verbose:
                print(f'Using existing file "{canvas_file}".')

        # Parse CSV
        dtype = [("name", "U50"), ("section", int), ("group", int)]
        conv = {0: _format_name,                 # name
                4: lambda name: name[-3:] or 0,  # section
                5: lambda name: name[-1:] or 0}  # group_name
        #cols = list(sorted(conv.keys()))
        names, sect, groups = np.loadtxt(canvas_file,
                                         delimiter=",", quotechar='"',
                                         dtype=dtype, converters=conv,
                                         skiprows=1, usecols=conv.keys(),
                                         unpack=True)

        for section in sections:
            mask = sect == section
            fig = _draw(names[mask], groups[mask],
                    title=f"Groups for Lab {lab:02d} Section {section:03d}")
            for ext in extensions:
                filename = os.path.join(f"lab{lab:02d}",
                                        f"groups{section:03d}.{ext}")
                fig.savefig(filename)
                if verbose:
                    print(f'Output written to "{filename}"')
            plt.close(fig)

sheets.parser = _sheets_parser


def _introduction_parser(parser):
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="print status messages")
    parser.add_argument("-u", "--update", action="store_true",
                        help="update quiz code")
    parser.add_argument("-l", "--lab", type=int,
                        default=max(existing_labs, default=0),
                        metavar="number", help="the lab's number")
    parser.add_argument("-s", "--sections", type=int, nargs="+",
                        default=[15, 25], metavar="section",
                        help="your section numbers")


def introduction(lab, sections, update=False, verbose=False):
    """Create a template for introduction slides

    The template has three slides:
        - title page with lab and first section number
        - sign-in sheets stacked on top of each other
        - a placeholder for the quiz code

    Since the quiz code changes quite frequently we put a placeholder in the
    template. Use the quiz_code command to update it before class.
    """
    intros_path = r"C:\\Users\\umthr\\OneDrive - Umich\\Documents\\Teaching" \
                  r"\\WN26 PHYSICS 251\\Introductions"
    template = intros_path + r"\\Template.pptx"

    if verbose:
        print(f'Using template at "{template}".')
    # Load template presentation
    prs = Presentation(template)

    # Modify title slide
    title_slide = prs.slides[0]
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Lab {lab:02d} - Section {section:03d}"

    # Modify group slide
    group_slide = prs.slides[1]
    try:  # looking for an existing picture shape
        pic_shape = next((shape for shape in group_slide.shapes
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
        left   = pic_shape.left
        top    = pic_shape.top
        width  = pic_shape.width
        height = pic_shape.height
        group_slide.shapes._spTree.remove(pic_shape._element)
        if verbose:
            print("Replacing existing picture on group slide.")
    except StopIteration:
        left   = 2487705
        top    = 0
        width  = 6656295
        height = 5143500
        if verbose:
            print("No existing picture found on group slide. "
                  "Falling back to default position and size.")

    # Add sign-in sheets for all sections on top of each other (first on top)
    for section in reversed(sections):
        img_path = f"lab{lab:02d}\\groups{section:03d}.png"
        if not os.path.exists(img_path):
            sheets([lab], [section], extensions=["png"])
        group_slide.shapes.add_picture(img_path, left, top,
                                       width=width, height=height)
        
    if update:  # quiz code
        quiz_slide = prs.slides[2]
        quiz_code = _get_quiz_code(lab, verbose=verbose)
        quiz_slide.placeholders[0].text = quiz_code

    # Save the modified presentation
    prs.save(f"{intros_path}\\PHYS251 Lab {lab:02d}.pptx")

introduction.parser = _introduction_parser


def _get_quiz_code(lab, verbose=False):
    # The quiz code for each lab is defined in a quiz on Canvas
    quizzes = _canvas_api(f"courses/{COURSE_ID}/quizzes",
                    headers={"search_term": f"Quiz {lab:d}:"},
                    verbose=verbose)
    try:
        quiz = next(quiz for quiz in quizzes
                    if quiz["title"].startswith(f"Quiz {lab:d}:"))
    except StopIteration:
        raise RuntimeError(f"Couldn't find quiz for lab {lab:d} on Canvas.")
    return quiz["access_code"]


def _quiz_code_parser(parser):
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="print status messages")
    parser.add_argument("-l", "--lab", type=int,
                        default=max(existing_labs, default=0),
                        metavar="number", help="the lab's number")


def quiz_code(lab, verbose=False):
    """Update the quiz code on the introduction slides
    
    This commands pulls the latest quiz code from the Canvas API and updates
    the corresponding slide in the introduction.
    """
    if verbose:
        print(f"Updating quiz code for lab {lab:d}.")

    intros_path = r"C:\\Users\\umthr\\OneDrive - Umich\\Documents\\Teaching" \
                  r"\\WN26 PHYSICS 251\\Introductions"
    intro = f"{intros_path}\\PHYS251 Lab {lab:02d}.pptx"

    if not os.path.exists(intro):
        raise RuntimeError(f'No slides for lab {lab:02d} found at "{intro}".')

    # Load template presentation
    prs = Presentation(intro)

    # Update quiz slide
    quiz_slide = prs.slides[2]
    quiz_code = _get_quiz_code(lab, verbose=verbose)
    quiz_slide.placeholders[0].text = quiz_code
    if verbose:
        print(f'Quiz code "{quiz_code}" retrieved from Canvas.')

    # Save the modified presentation
    prs.save(intro)
    if verbose:
        print(f'Updated presentation "{intro}".')

quiz_code.parser = _quiz_code_parser

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Utilities for using Canvas "
                                                 "as a GSI")
    
    commands = [sheets, introduction, quiz_code]
    aliases = {"introduction": ["intro"],
               "quiz_code": ["quiz"]}
    subparsers = parser.add_subparsers(dest="command",
                                       title="commands",
                                       required=True)
    for command in commands:
        description, epilog = command.__doc__.split("\n\n", maxsplit=1)
        subparser = subparsers.add_parser(command.__name__,
                                          aliases=aliases.get(command.__name__,
                                                              []),
                                          formatter_class=
                                          RawDescriptionDefaultsHelpFormatter,
                                          description=description,
                                          help=description,
                                          epilog=epilog)
        # Populate subparser with command-specific arguments 
        # see e.g., sheets.parser = _sheets_parser
        command.parser(subparser)
    
    args = vars(parser.parse_args())

    command_name = args.pop("command")
    try:
        command = next(command for command in commands
                       if command.__name__ == command_name)
    except StopIteration:
        raise RuntimeError(f'Unknown command: "{command_name}"')
    
    command(**args)